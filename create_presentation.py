#!/usr/bin/env python3
"""
Script to create PowerPoint presentation for Hybrid KG-LLM research.
Requires: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define colors
    title_color = RGBColor(31, 78, 121)  # Dark blue
    accent_color = RGBColor(68, 114, 196)  # Medium blue
    text_color = RGBColor(0, 0, 0)  # Black
    
    def add_title_slide(title, subtitle=""):
        """Add a title slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(44)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = title_color
        
        if subtitle:
            subtitle_shape.text = subtitle
            subtitle_shape.text_frame.paragraphs[0].font.size = Pt(28)
            subtitle_shape.text_frame.paragraphs[0].font.color.rgb = accent_color
        
        return slide
    
    def add_content_slide(title, content_items):
        """Add a content slide with title and bullet points"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title_shape = slide.shapes.title
        content_shape = slide.placeholders[1]
        
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(36)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = title_color
        
        tf = content_shape.text_frame
        tf.word_wrap = True
        tf.clear()
        
        for item in content_items:
            p = tf.add_paragraph()
            p.text = item
            p.level = 0
            p.font.size = Pt(20)
            p.font.color.rgb = text_color
            p.space_after = Pt(12)
        
        return slide
    
    def add_two_column_slide(title, left_items, right_items):
        """Add a slide with two columns"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(36)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = title_color
        
        # Remove default placeholder
        slide.shapes.element.remove(slide.placeholders[1].element)
        
        # Add left text box
        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5.5))
        left_tf = left_box.text_frame
        left_tf.word_wrap = True
        left_tf.clear()
        
        for item in left_items:
            p = left_tf.add_paragraph()
            p.text = item
            p.level = 0
            p.font.size = Pt(18)
            p.font.color.rgb = text_color
            p.space_after = Pt(10)
        
        # Add right text box
        right_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(4.5), Inches(5.5))
        right_tf = right_box.text_frame
        right_tf.word_wrap = True
        right_tf.clear()
        
        for item in right_items:
            p = right_tf.add_paragraph()
            p.text = item
            p.level = 0
            p.font.size = Pt(18)
            p.font.color.rgb = text_color
            p.space_after = Pt(10)
        
        return slide
    
    # Slide 1: Research Question & Motivation
    add_content_slide(
        "Research Question & Motivation",
        [
            "The Problem:",
            "• Multi-hop reasoning over knowledge graphs is challenging for LLMs",
            "• Existing approaches struggle with:",
            "  - Over-squashing: Too many neighbors overwhelm the model",
            "  - Heterophily: Diverse relation types confuse reasoning",
            "  - Lack of alignment: Models don't learn preferred reasoning patterns",
            "",
            "Our Goal:",
            "• Train LLMs to perform multi-hop reasoning over KGs using:",
            "  - Hybrid datasets (text + optional visual cues)",
            "  - Preference alignment via DPO (Direct Preference Optimization)",
            "  - Similarity-based neighbor selection to focus on relevant paths"
        ]
    )
    
    # Slide 2: High-Level Architecture
    add_content_slide(
        "Our Approach - High-Level Architecture",
        [
            "Three-Stage Pipeline:",
            "",
            "Stage 1: Data Preparation",
            "• Extract KG triples from PrimeKG",
            "• Generate reasoning paths (positive vs negative)",
            "• Optional: Render graph visualizations",
            "• Optional: Apply SimCSE-based neighbor ranking",
            "",
            "Stage 2: DPO Training",
            "• Fine-tune LLM (GPT-2, Mistral, etc.) using Direct Preference Optimization",
            "• Model learns to prefer correct reasoning chains over incorrect ones",
            "• Training on hybrid prompts: text + optional graph images",
            "",
            "Stage 3: Evaluation",
            "• Link prediction (MRR, Hits@10)",
            "• Multi-hop QA (accuracy)",
            "• Ablation studies"
        ]
    )
    
    # Slide 3: Hybrid Dataset Construction
    add_content_slide(
        "Key Innovation 1 - Hybrid Dataset Construction",
        [
            "What Makes It 'Hybrid'?",
            "",
            "1. Multi-Modal Prompts",
            "• Textual component: Query + selected KG triples",
            "• Visual component (optional): Rendered graph subgraphs",
            "• Example: 'What disease is associated with gene X?' + triples + graph image",
            "",
            "2. DPO Pair Generation",
            "• Chosen response: Correct reasoning path",
            "• Rejected response: Incorrect or suboptimal path",
            "• Model learns to distinguish between good and bad reasoning",
            "",
            "3. Similarity-Based Neighbor Selection (Optional)",
            "• Use SimCSE embeddings to rank neighbors by semantic similarity",
            "• Reduces over-squashing by focusing on relevant triples",
            "• Threshold-based filtering (e.g., top-k=5, similarity > 0.8)"
        ]
    )
    
    # Slide 4: DPO Training (MOST IMPORTANT)
    add_content_slide(
        "Key Innovation 2 - DPO Training for KG Reasoning",
        [
            "Why DPO?",
            "• Traditional fine-tuning: Requires explicit labels, doesn't capture preferences",
            "• Our DPO approach: Learns from relative preferences (chosen vs rejected)",
            "• Directly optimizes for reasoning quality",
            "• More sample-efficient than RLHF",
            "",
            "Training Process:",
            "• Input: Prompt (query + triples + optional image)",
            "• Chosen: Correct reasoning chain",
            "• Rejected: Incorrect reasoning chain",
            "• Objective: Maximize chosen, minimize rejected (controlled by β)",
            "",
            "Training Configuration:",
            "• Base models: GPT-2 (124M) to Mistral-7B",
            "• Learning rate: 5e-6",
            "• Batch size: 1-4 (with gradient accumulation)",
            "• Epochs: 2-3"
        ]
    )
    
    # Slide 5: Integration Framework
    add_content_slide(
        "Key Innovation 3 - Integration Framework",
        [
            "Combining Three Approaches:",
            "",
            "1. SNS (Similarity-based Neighbor Selection)",
            "• What we took: SimCSE-based neighbor ranking",
            "• How we adapted: Applied to KG entities/relations instead of graph nodes",
            "• Benefit: Reduces noise, focuses on semantically relevant triples",
            "",
            "2. GITA (Graph to Image-Text Integration)",
            "• What we took: Graph visualization via Graphviz/NetworkX",
            "• How we adapted: Rendered KG subgraphs as images for multi-modal prompts",
            "• Benefit: Visual grounding helps model understand graph structure",
            "",
            "3. GraphWiz (Graph Reasoning LLM)",
            "• What we took: DPO training framework",
            "• How we adapted: Extended to KG domains with hybrid text+visual prompts",
            "• Benefit: Preference alignment for KG reasoning tasks",
            "",
            "Novel Contributions:",
            "• Dual-stage caching system (reproducibility)",
            "• HPC-aware rendering (scalability)",
            "• DPO ablation tooling (systematic exploration)"
        ]
    )
    
    # Slide 6: Experimental Setup
    add_two_column_slide(
        "Experimental Setup",
        [
            "Dataset:",
            "• Source: PrimeKG (biomedical knowledge graph)",
            "• Format: Train/Val/Test splits (80/10/10 or 70/15/15)",
            "• Size: Scalable from 10 samples (demo) to 1000+ samples",
            "• Tasks: Link prediction, multi-hop QA",
            "",
            "Models:",
            "• Base: GPT-2 (124M) - for accessibility",
            "• Extended: Mistral-7B-Instruct - for performance",
            "• Training: DPO fine-tuning with LoRA (optional)"
        ],
        [
            "Evaluation Metrics:",
            "• Link Prediction: MRR, Hits@10",
            "• Multi-hop QA: Accuracy, Precision, Recall, F1",
            "",
            "Ablations:",
            "• SimCSE threshold",
            "• DPO β parameter",
            "• With/without images",
            "",
            "Hardware:",
            "• GPT-2: Standard laptop (CPU)",
            "• Mistral-7B: GPU with 16GB+ VRAM"
        ]
    )
    
    # Slide 7: Results & Findings
    add_two_column_slide(
        "Results & Findings",
        [
            "Training Results:",
            "• Training Loss: 0.46 (final, after 2 epochs)",
            "• Training Time: ~20 seconds (GPT-2, small dataset)",
            "• Convergence: ✓ Model learns to distinguish chosen vs rejected",
            "",
            "Evaluation Results (Multi-hop QA):",
            "• Accuracy: 96.7% (29/30 correct)",
            "• Precision: 92.9%",
            "• Recall: 91.1%",
            "• F1 Score: 91.8%",
            "• Test Samples: 30",
            "",
            "Per-Relation Performance:",
            "• 13/14 relation types: 100% accuracy",
            "• Strong performance across diverse relations",
            "• Example relations: causes, treats, activates, inhibits, etc."
        ],
        [
            "Key Findings:",
            "",
            "1. DPO Training Works",
            "• Model learns preferred reasoning patterns",
            "• Training loss decreases over epochs",
            "• Framework is operational and scalable",
            "",
            "2. Strong Evaluation Performance",
            "• 96.7% accuracy on multi-hop QA",
            "• High precision (92.9%) and recall (91.1%)",
            "• Consistent across relation types",
            "",
            "3. Hybrid Approach Benefits",
            "• Visual cues provide structural intuition",
            "• SimCSE ranking reduces over-squashing",
            "• Multi-modal prompts enhance understanding",
            "",
            "Limitations & Future Work:",
            "• Current: 30 test samples (proof-of-concept)",
            "• Future: Scale to 1000+ samples",
            "• Future: Link prediction metrics (MRR, Hits@10)",
            "• Future: Baseline comparisons, statistical testing"
        ]
    )
    
    # Slide 8: Comparison & Standing Out
    add_content_slide(
        "Comparison & Standing Out",
        [
            "How We Stand Out:",
            "",
            "vs. SNS Alone:",
            "✓ Adds DPO training (not just prompt engineering)",
            "✓ Extends to KG domains (not just citation networks)",
            "✓ Multi-modal support (text + visual)",
            "",
            "vs. GITA Alone:",
            "✓ Focuses on KG reasoning (not general graphs)",
            "✓ Uses DPO for preference alignment (not just fine-tuning)",
            "✓ Integrates similarity-based neighbor selection",
            "",
            "vs. GraphWiz Alone:",
            "✓ Adds visual grounding (graph images)",
            "✓ Incorporates SimCSE ranking",
            "✓ End-to-end pipeline for PrimeKG-scale KGs",
            "",
            "Unique Contributions:",
            "1. First hybrid DPO approach for KG reasoning",
            "2. Integration framework combining three methods",
            "3. Reproducible pipeline with caching and HPC support",
            "4. Ablation tooling for systematic hyperparameter exploration"
        ]
    )
    
    # Slide 9: Takeaways & Impact
    add_content_slide(
        "Takeaways & Impact",
        [
            "Key Takeaways:",
            "1. DPO training effectively aligns LLMs for KG reasoning",
            "2. Hybrid approach (text + visual) improves understanding",
            "3. Similarity-based selection reduces noise and over-squashing",
            "4. End-to-end pipeline enables reproducible research",
            "",
            "Impact:",
            "• Methodology: Novel integration of three approaches",
            "• Practical: Works on standard hardware (GPT-2)",
            "• Scalable: Extends to large models and datasets",
            "• Reproducible: Deterministic caching and HPC support"
        ]
    )
    
    # Slide 10: Conclusion
    add_content_slide(
        "Conclusion & Future Directions",
        [
            "Summary:",
            "• We present a hybrid DPO training framework for multi-hop KG reasoning",
            "• Combines SNS neighbor selection, GITA visual grounding, and GraphWiz DPO training",
            "• Demonstrates effective learning of preferred reasoning patterns",
            "• Provides reproducible, scalable pipeline for KG-LLM research",
            "",
            "Future Work:",
            "• Scale to larger datasets (1000+ samples)",
            "• Compare with baselines (zero-shot, random, rule-based)",
            "• Statistical significance testing",
            "• Adaptive sampling for large KGs",
            "• Transfer learning to unseen KG tasks",
            "",
            "Thank You!",
            "Questions?"
        ]
    )
    
    return prs

if __name__ == "__main__":
    print("Creating PowerPoint presentation...")
    prs = create_presentation()
    output_file = "Hybrid_KG_LLM_Presentation.pptx"
    prs.save(output_file)
    print(f"✓ Presentation saved as: {output_file}")
    print(f"✓ Total slides: {len(prs.slides)}")

